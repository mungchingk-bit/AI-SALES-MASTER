import csv
import json
import logging
import os
import re

from models.chat_message import ChatMessage


def parse_file(file_path: str) -> list[ChatMessage]:
    """Parse a sales chat file into a list of ChatMessage. Supports txt/csv/json/docx/pdf/pptx."""
    ext = file_path.lower().split(".")[-1]

    if ext == "json":
        return _parse_json(file_path)
    elif ext == "csv":
        return _parse_csv(file_path)
    elif ext in ("docx", "doc"):
        return _parse_docx(file_path)
    elif ext == "pdf":
        return _parse_pdf(file_path)
    elif ext in ("pptx", "ppt"):
        return _parse_pptx(file_path)
    else:
        return _parse_txt(file_path)


def extract_text(file_path: str) -> str:
    """Extract raw text from any supported file. For knowledge base ingestion."""
    ext = file_path.lower().split(".")[-1]

    if ext in ("docx", "doc"):
        return _extract_docx_text(file_path)
    elif ext == "pdf":
        return _extract_pdf_text(file_path)
    elif ext in ("pptx", "ppt"):
        return _extract_pptx_text(file_path)
    elif ext in ("xlsx", "xls"):
        return _extract_xlsx_text(file_path)
    elif ext in ("jpg", "jpeg", "png"):
        return _extract_image_text(file_path)
    else:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


def _parse_txt(file_path: str) -> list[ChatMessage]:
    """Parse a text file with '销售：/客户：' prefixed lines."""
    messages = []
    with open(file_path, "r", encoding="utf-8") as f:
        lines = f.readlines()

    current_speaker = None
    current_content = []

    for line in lines:
        line = line.strip()
        if not line:
            continue

        detected_speaker = _detect_speaker(line)
        if detected_speaker:
            if current_speaker and current_content:
                messages.append(
                    ChatMessage(role=_map_role(current_speaker), content="".join(current_content).strip())
                )
            current_speaker = detected_speaker
            content = _remove_speaker_prefix(line, detected_speaker)
            current_content = [content] if content else []
        else:
            if current_speaker:
                current_content.append(line)
            else:
                if len(messages) % 2 == 0:
                    current_speaker = "sales"
                else:
                    current_speaker = "customer"
                current_content = [line]

    if current_speaker and current_content:
        messages.append(
            ChatMessage(role=_map_role(current_speaker), content="".join(current_content).strip())
        )

    return messages


def _parse_csv(file_path: str) -> list[ChatMessage]:
    messages = []
    with open(file_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        for row in reader:
            speaker = row.get("speaker", "").strip()
            content = row.get("content", "").strip()
            if speaker and content:
                messages.append(ChatMessage(role=_map_role(speaker), content=content))
    return messages


def _parse_json(file_path: str) -> list[ChatMessage]:
    messages = []
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    if isinstance(data, list):
        for item in data:
            role = item.get("role", "")
            content = item.get("content", "")
            if role and content:
                messages.append(ChatMessage(role=_map_role(role), content=content))
    return messages


def _parse_docx(file_path: str) -> list[ChatMessage]:
    """Parse docx as conversation. Tries to detect speaker patterns first,
    falls back to treating as sales monologue/script."""
    text = _extract_docx_text(file_path)
    return _parse_text_as_conversation(text)


def _parse_pdf(file_path: str) -> list[ChatMessage]:
    """Parse PDF as conversation."""
    text = _extract_pdf_text(file_path)
    return _parse_text_as_conversation(text)


def _parse_pptx(file_path: str) -> list[ChatMessage]:
    """Parse PPT as conversation."""
    text = _extract_pptx_text(file_path)
    return _parse_text_as_conversation(text)


def _parse_text_as_conversation(text: str) -> list[ChatMessage]:
    """Parse extracted text into conversation messages.
    Handles multiple formats: speaker-prefixed, Q&A, dialogue turns."""
    messages = []

    # Try to detect "说话人1："/"说话人2：" patterns (from voice transcription)
    speaker_pattern = re.compile(r"(说话人\d|销售[：:]|客户[：:]|问[：:]|答[：:]|Q[：:]|A[：:])")
    lines = text.split("\n")

    has_speakers = any(speaker_pattern.search(line) for line in lines[:50])

    if has_speakers:
        current_speaker = None
        current_content = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            match = speaker_pattern.match(line)
            if match:
                if current_speaker and current_content:
                    messages.append(ChatMessage(
                        role=_map_speaker_pattern(current_speaker),
                        content=" ".join(current_content).strip()
                    ))
                current_speaker = match.group(1)
                content = speaker_pattern.sub("", line).strip()
                current_content = [content] if content else []
            else:
                if current_speaker:
                    current_content.append(line)
                elif len(messages) % 2 == 0:
                    current_speaker = "销售："
                    current_content = [line]
                else:
                    current_speaker = "客户："
                    current_content = [line]

        if current_speaker and current_content:
            messages.append(ChatMessage(
                role=_map_speaker_pattern(current_speaker),
                content=" ".join(current_content).strip()
            ))
    else:
        # No speaker patterns - treat as a single sales script
        # Split by paragraphs
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        for p in paragraphs:
            messages.append(ChatMessage(role="user", content=p))

    return messages


def _map_speaker_pattern(speaker: str) -> str:
    """Map voice transcription speaker labels to user/assistant."""
    s = speaker.rstrip("：:")
    if s in ("销售", "答", "A"):
        return "user"
    elif s in ("客户", "问", "Q"):
        return "assistant"
    # 说话人1/2 etc - alternate
    num = re.search(r"\d", s)
    if num:
        return "user" if int(num.group()) % 2 == 1 else "assistant"
    return "user"


def _extract_docx_text(file_path: str) -> str:
    """Extract text from a docx file."""
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    except ImportError:
        # Fallback: use zipfile to extract raw XML text
        import zipfile
        text_parts = []
        with zipfile.ZipFile(file_path) as z:
            for name in z.namelist():
                if name.endswith(".xml") and "word/" in name:
                    content = z.read(name).decode("utf-8", errors="replace")
                    # Strip XML tags
                    text = re.sub(r"<[^>]+>", "", content)
                    text = re.sub(r"\s+", " ", text).strip()
                    if text:
                        text_parts.append(text)
        return "\n".join(text_parts)


def _extract_pdf_text(file_path: str) -> str:
    """Extract text from a PDF file."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        text_parts = []
        for page in doc:
            text_parts.append(page.get_text())
        return "\n".join(text_parts)
    except ImportError:
        try:
            from pdfminer.high_level import extract_text
            return extract_text(file_path)
        except ImportError:
            return f"[无法解析PDF，请安装PyMuPDF: pip install PyMuPDF] 文件：{os.path.basename(file_path)}"


def _extract_pptx_text(file_path: str) -> str:
    """Extract text from a PPT/PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                text_parts.append(f"--- 第{i}页 ---\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        return f"[无法解析PPT，请安装python-pptx: pip install python-pptx] 文件：{os.path.basename(file_path)}"


def _detect_speaker(line: str) -> str | None:
    prefixes = {
        "销售": "sales", "我方": "sales", "Sales": "sales",
        "客户": "customer", "对方": "customer", "Customer": "customer",
    }
    for cn_prefix, speaker in prefixes.items():
        if line.startswith(cn_prefix + "：") or line.startswith(cn_prefix + ":"):
            return speaker
    return None


def _remove_speaker_prefix(line: str, speaker: str) -> str:
    prefix_map = {
        "sales": ["销售：", "销售:", "我方：", "我方:", "Sales:", "Sales："],
        "customer": ["客户：", "客户:", "对方：", "对方:", "Customer:", "Customer："],
    }
    for prefix in prefix_map.get(speaker, []):
        if line.startswith(prefix):
            return line[len(prefix):].strip()
    return line


def _map_role(speaker: str) -> str:
    sales_ids = {"sales", "销售", "我方", "sale"}
    customer_ids = {"customer", "客户", "对方", "client"}
    s = speaker.lower().strip()
    if s in sales_ids:
        return "user"
    elif s in customer_ids:
        return "assistant"
    return "user"


def _extract_xlsx_text(file_path: str) -> str:
    """Extract text from an XLSX/XLS file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_lines = [f"--- Sheet: {sheet_name} ---"]
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                line = " | ".join(cells)
                if line.strip(" |"):
                    sheet_lines.append(line)
            text_parts.append("\n".join(sheet_lines))
        wb.close()
        return "\n\n".join(text_parts)
    except ImportError:
        return f"[无法解析XLSX，请安装openpyxl: pip install openpyxl] 文件：{os.path.basename(file_path)}"


def _extract_image_text(file_path: str) -> str:
    """Extract text from an image file using OCR."""
    import config
    if config.OCR_MODE == "tesseract":
        return _ocr_tesseract(file_path)
    else:
        return _ocr_ollama_vision(file_path)


def _ocr_tesseract(file_path: str) -> str:
    """Use pytesseract for OCR."""
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return text.strip() if text.strip() else "[图片中未检测到文字]"
    except ImportError:
        return "[无法进行OCR，请安装pytesseract: pip install pytesseract，并安装Tesseract OCR]"
    except Exception as e:
        return f"[OCR处理失败: {str(e)}]"


def _ocr_ollama_vision(file_path: str) -> str:
    """Use Ollama vision model for OCR."""
    import base64
    import requests
    import config

    with open(file_path, "rb") as f:
        image_b64 = base64.b64encode(f.read()).decode("utf-8")

    ext = file_path.lower().split(".")[-1]
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png"}
    mime_type = mime_map.get(ext, "image/jpeg")

    payload = {
        "model": config.OLLAMA_VISION_MODEL,
        "messages": [
            {
                "role": "user",
                "content": "请仔细识别并提取这张图片中的所有文字内容，按原排版输出。如果图片中没有文字，描述图片内容。",
                "images": [image_b64],
            }
        ],
        "stream": False,
    }

    try:
        response = requests.post(
            f"{config.OLLAMA_BASE_URL}/api/chat",
            json=payload,
            timeout=120,
        )
        response.raise_for_status()
        result = response.json()["message"]["content"]
        return result.strip() if result.strip() else "[图片中未检测到文字]"
    except Exception as e:
        logging.warning(f"Ollama vision OCR failed: {e}, falling back to tesseract")
        return _ocr_tesseract(file_path)
